"""Low-level primitives over python-pptx.

These helpers are layout-agnostic and brand-agnostic. They cover the boring
boilerplate (filling shapes, anchoring text frames, paragraph spacing) so that
layout modules can read like declarative compositions.

Standard slide size used everywhere : 16:9 widescreen
    339.9 mm × 190.5 mm  (= 33.867 cm × 19.05 cm)
"""
from __future__ import annotations

import os
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Standard slide dimensions (in cm)
SLIDE_W_CM = 33.867
SLIDE_H_CM = 19.05

# ---------------------------------------------------------------------------
# Shape fills / lines

def set_fill(shape, rgb: RGBColor):
    """Solid fill, no border, flat (no inherited theme shadow — kit is flat)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    shape.shadow.inherit = False


def set_line(shape, rgb: RGBColor, width_pt: float = 0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x_cm, y_cm, w_cm, h_cm, *, fill: RGBColor | None = None):
    """Add a filled rectangle. Returns the shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm)
    )
    if fill is not None:
        set_fill(shape, fill)
    return shape


def add_round_rect(slide, x_cm, y_cm, w_cm, h_cm, *,
                   fill: RGBColor | None = None, radius_frac: float = 0.12):
    """Add a rounded rectangle. `radius_frac` is the corner radius as a fraction
    of the shorter side (0.5 = fully rounded pill). Returns the shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm)
    )
    try:
        shape.adjustments[0] = max(0.0, min(0.5, radius_frac))
    except (IndexError, ValueError):
        pass
    if fill is not None:
        set_fill(shape, fill)
    return shape


def add_oval(slide, x_cm, y_cm, d_cm, *, fill: RGBColor | None = None,
             line: RGBColor | None = None):
    """Add a small circle/disc of diameter `d_cm`. Used for rating indicators."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(x_cm), Cm(y_cm), Cm(d_cm), Cm(d_cm)
    )
    if fill is not None:
        set_fill(shape, fill)
    if line is not None:
        set_line(shape, line, width_pt=1.0)
        shape.fill.background()
        shape.shadow.inherit = False
    return shape


def set_hanging(p, marL_cm: float):
    """Hanging indent : wrapped lines align under the text, not the marker.
    `marL_cm` = left margin of the text column (marker sits in the hang)."""
    from pptx.util import Cm as _Cm
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    pPr.set("marL", str(int(_Cm(marL_cm))))
    pPr.set("indent", str(-int(_Cm(marL_cm))))


# ---------------------------------------------------------------------------
# Text frames

ANCHOR = {
    "t": MSO_ANCHOR.TOP,
    "m": MSO_ANCHOR.MIDDLE,
    "b": MSO_ANCHOR.BOTTOM,
}

ALIGN = {
    "left":   PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right":  PP_ALIGN.RIGHT,
}


def add_text(slide, x_cm, y_cm, w_cm, h_cm, *,
             anchor: str = "t",
             margins: tuple[float, float, float, float] = (0.1, 0.15, 0.1, 0.15)):
    """Insert a text box. Returns (shape, text_frame).

    margins : (left, top, right, bottom) in cm.
    anchor  : "t" (top, default), "m" (middle), or "b" (bottom).
    """
    box = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tf = box.text_frame
    tf.word_wrap = True
    ml, mt, mr, mb = margins
    tf.margin_left, tf.margin_top, tf.margin_right, tf.margin_bottom = (
        Cm(ml), Cm(mt), Cm(mr), Cm(mb)
    )
    tf.vertical_anchor = ANCHOR[anchor]
    return box, tf


def para(tf, *, first: bool = False, align: str | None = None,
         space_before: float = 0, space_after: float = 2,
         line_spacing: float = 1.15):
    """Get or add a paragraph in a text frame."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align:
        p.alignment = ALIGN[align]
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    p.line_spacing = line_spacing
    return p


def run(p, text: str, *,
        size: float = 10,
        color: RGBColor,
        bold: bool = False,
        italic: bool = False,
        font: str):
    """Append a styled run to a paragraph.

    `color` and `font` are required — they should come from the active Charte.
    """
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


# ---------------------------------------------------------------------------
# Images

def add_image(slide, path: str | os.PathLike, x_cm: float, y_cm: float, *,
              w_cm: float | None = None, h_cm: float | None = None):
    """Insert an image. If only one of (w_cm, h_cm) is given, aspect ratio is
    preserved. Returns None if the file is missing."""
    if not os.path.exists(path):
        return None
    kwargs = {}
    if w_cm is not None: kwargs["width"]  = Cm(w_cm)
    if h_cm is not None: kwargs["height"] = Cm(h_cm)
    return slide.shapes.add_picture(str(path), Cm(x_cm), Cm(y_cm), **kwargs)


def fit_image(pic, max_w_cm: float | None = None, max_h_cm: float | None = None):
    """Shrink an inserted picture to fit within a bounding box, preserving aspect.
    Safe no-op if the picture is already smaller."""
    if pic is None:
        return
    if max_w_cm is not None and pic.width > Cm(max_w_cm):
        ratio = pic.width / pic.height
        pic.width  = Cm(max_w_cm)
        pic.height = int(pic.width / ratio)
    if max_h_cm is not None and pic.height > Cm(max_h_cm):
        ratio = pic.width / pic.height
        pic.height = Cm(max_h_cm)
        pic.width  = int(pic.height * ratio)


def right_align(pic, slide_right_cm: float, margin_cm: float = 0):
    """Reposition an inserted picture so its right edge sits at `slide_right_cm`."""
    if pic is None:
        return
    pic.left = Cm(slide_right_cm - margin_cm) - pic.width


# ---------------------------------------------------------------------------
# Slide creation

def add_blank_slide(prs):
    """Add a slide using the blank layout (last in default layouts)."""
    return prs.slides.add_slide(prs.slide_layouts[6])
