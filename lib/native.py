"""NativeDeck — construit un deck en remplissant les layouts d'un thème INGÉRÉ.

Pendant « build » d'`ingest` (lib/template.py). Là où `Deck` dessine des formes
sur un PPTX vide via les layouts Python du kit, `NativeDeck` :
  • charge le `template.pptx` du thème (master + slideLayouts du client préservés) ;
  • vide les slides d'exemple ;
  • instancie les layouts voulus et remplit leurs placeholders PAR RÔLE
    (title / eyebrow / subtitle / content / picture / footer / date), via le
    rôle déduit par position — jamais par idx (robuste aux variantes +IMG) ;
  • `fit_text` (autofit) pour ne jamais déborder.

C'est la généralisation de l'ancien `build_native.py` (codé en dur pour Egis) :
ici rien n'est spécifique à un template, tout vient du thème + du catalogue.

    deck = NativeDeck("egis")
    deck.add(deck.pick(kind="cover"), title="…", eyebrow="…", date="…")
    deck.add("Chapter A [Dark pict]", eyebrow="Partie 01", title="…",
             number="01", photo="…/chap1.jpg")
    deck.add(deck.pick(kind="grid", min_content=4), title="…", subtitle="…",
             cells=[("T","desc"), …], icons=[…])
    deck.add(deck.pick(kind="list"), title="…", items=[("nom","desc"), …])
    deck.save("out/deck.pptx"); deck.export_pdf("out/deck.pptx")
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_AUTO_SIZE

from lib.charte import _theme_search_paths
from lib.template import _role

_PIC = "PICTURE (18)"


def _find_theme(name: str) -> Path:
    for base in _theme_search_paths():
        d = base / name
        if (d / "template.pptx").exists():
            return d
    raise FileNotFoundError(
        f"thème natif '{name}' introuvable (cherché template.pptx dans : "
        + ", ".join(str(b / name) for b in _theme_search_paths()) + ")")


def _has_image(ph) -> bool:
    return ph._element.find(".//" + qn("a:blip")) is not None


class NativeDeck:
    def __init__(self, theme: str):
        self.dir = _find_theme(theme)
        self.catalog = json.loads((self.dir / "catalog.json").read_text("utf-8"))["layouts"]
        self.prs = Presentation(str(self.dir / "template.pptx"))
        self._layouts = {lay.name: lay
                         for m in self.prs.slide_masters for lay in m.slide_layouts}
        # vider les slides d'exemple (garder masters + layouts)
        lst = self.prs.slides._sldIdLst
        for sid in list(lst):
            self.prs.part.drop_rel(sid.get(qn("r:id")))
            lst.remove(sid)

    # ---- sélection de layout

    def info(self, name: str) -> dict | None:
        return next((l for l in self.catalog if l["name"] == name), None)

    def pick(self, *, kind: str | None = None, min_content: int | None = None,
             name: str | None = None) -> str | None:
        """Choisit un layout par nom exact, ou par kind/capacité (1er match du catalogue)."""
        if name:
            return name if name in self._layouts else None
        for l in self.catalog:
            if kind and not l["kind"].startswith(kind):
                continue
            if min_content is not None and l["signature"]["content"] < min_content:
                continue
            return l["name"]
        return None

    # ---- rendu

    def _roles(self, slide) -> dict[str, list]:
        title = next((p for p in slide.placeholders
                      if str(p.placeholder_format.type) in ("TITLE (1)", "CENTER_TITLE (3)")), None)
        tt = title.top if title is not None else None
        tl = title.left if title is not None else None
        out: dict[str, list] = {}
        for ph in slide.placeholders:
            out.setdefault(_role(ph, tt, tl), []).append(ph)
        for lst in out.values():
            lst.sort(key=lambda p: (round((p.top or 0) / 100000), p.left or 0))
        return out

    @staticmethod
    def _fit(tf):
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def _txt(self, ph, text):
        tf = ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
        self._fit(tf)

    def _levels(self, ph, entries):
        tf = ph.text_frame
        tf.clear()
        for i, (lvl, t) in enumerate(entries):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = t
            p.level = lvl
        self._fit(tf)

    def _icon(self, slide, ph, icon):
        l, t, w, h = ph.left, ph.top, ph.width, ph.height
        ph._element.getparent().remove(ph._element)
        sz = int(min(w, h) * 0.82)
        slide.shapes.add_picture(icon, int(l + (w - sz) / 2), int(t + (h - sz) / 2), height=sz)

    def _target(self, slide, layout, role):
        """Placeholder de `role` sur la slide ; sinon MATÉRIALISÉ depuis le layout
        (python-pptx ne clone pas date/footer/page_number par défaut) ; sinon None."""
        got = self._roles(slide).get(role)
        if got:
            return got[0]
        lay = self._layouts[layout]
        for ph in self._roles(lay).get(role, []):
            try:
                slide.shapes.clone_placeholder(ph)
            except Exception:
                return None
            again = self._roles(slide).get(role)
            return again[0] if again else None
        return None

    def add(self, layout: str, *, title=None, eyebrow=None, subtitle=None,
            cells=None, items=None, number=None, footer=None, date=None,
            photo=None, icons=None):
        """Ajoute une slide depuis `layout` et remplit ses placeholders par rôle.

        cells  : [(titre, desc), …]  → une cellule par placeholder `content` (colonnes/blocs)
        items  : [(nom, desc), …]    → une LISTE dans le plus grand placeholder `content`
        number : str                 → gros numéro de section
        photo  : chemin              → remplit le plus grand cadre image (recadre)
        icons  : [chemins]           → icônes centrées dans les petits cadres image

        Toute valeur fournie sans placeholder cible (rôle absent du layout) est
        signalée (pas de perte silencieuse — cf. politique « no silent fallback »).
        """
        if layout not in self._layouts:
            raise KeyError(f"layout '{layout}' absent du thème")
        slide = self.prs.slides.add_slide(self._layouts[layout])
        missing: list[str] = []

        for val, role in ((title, "title"), (eyebrow, "eyebrow"),
                          (subtitle, "subtitle"), (date, "date"), (footer, "footer")):
            if val is None:
                continue
            ph = self._target(slide, layout, role)
            if ph is None:
                missing.append(role)
            else:
                self._txt(ph, val)

        content = self._roles(slide).get("content", [])
        if number is not None:
            (self._txt(content[0], number) if content else missing.append("content"))
        elif items is not None:
            if content:
                big = max(content, key=lambda p: (p.width or 0) * (p.height or 0))
                entries = []
                for name, desc in items:
                    entries += [(1, name), (2, desc)]
                self._levels(big, entries)
            else:
                missing.append("content")
        elif cells is not None:
            if not content:
                missing.append("content")
            for cell, ph in zip(cells, content):
                if isinstance(cell, (tuple, list)):
                    self._levels(ph, [(0, cell[0]), (1, cell[1])])
                else:
                    self._txt(ph, cell)

        pics = list(self._roles(slide).get("picture", []))
        if photo:
            if pics:
                big = max(pics, key=lambda p: (p.width or 0) * (p.height or 0))
                big.insert_picture(photo)
                pics.remove(big)
            else:
                missing.append("photo")
        if icons:
            for icon, ph in zip(icons, pics):
                self._icon(slide, ph, icon)

        # nettoyage : retirer les placeholders VIDES qui afficheraient un prompt
        # en édition — cadres image sans blip, et zones `content` non remplies.
        roles = self._roles(slide)
        empties = set(id(p) for p in roles.get("content", [])
                      if not p.text_frame.text.strip())
        for ph in list(slide.placeholders):
            t = str(ph.placeholder_format.type)
            if t == _PIC and not _has_image(ph):
                ph._element.getparent().remove(ph._element)
            elif id(ph) in empties:
                ph._element.getparent().remove(ph._element)

        if missing:
            print(f"[NativeDeck] slide « {layout} » : champ(s) fourni(s) sans "
                  f"placeholder cible, ignoré(s) : {', '.join(missing)}")
        return slide

    # ---- sortie

    def save(self, path: str) -> str:
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return str(path)

    def export_pdf(self, pptx_path: str, pdf_path: str | None = None):
        from lib.pdf_export import pptx_to_pdf, soffice_path
        if soffice_path() is None:
            return None
        return pptx_to_pdf(pptx_path, pdf_path)
