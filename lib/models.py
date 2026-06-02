"""ModelDeck — construit un deck en CLONANT des slides-modèles designées.

Pour les thèmes en mode « models » (cf. lib/template.ingest) : le template n'a
pas de bibliothèque de masques, son design est dessiné sur les slides. On réutilise
ces slides comme gabarits : `clone(ref)` duplique une slide-modèle, `.set({...})`
réécrit son texte par ancre (le texte actuel d'une forme), le design reste intact.

    deck = ModelDeck("otomata")
    deck.catalog()                       # voir les modèles + leurs ancres
    deck.clone("le contexte").set({      # ref = nom | index | kind
        "№ 01 — LE CONSTAT": "№ 01 — LE DÉFI",
        "le contexte": "le problème",
        "la plupart des projets": "Trop de chantiers IA restent des démos.",
    })
    deck.save("out/deck.pptx")

Le clone réécrit garde polices/couleurs/positions du modèle. Les originaux servent
de source puis sont retirés à `save()` — le deck final ne contient que les clones.
"""
from __future__ import annotations

import copy
import json

from pptx import Presentation
from pptx.oxml.ns import qn

from lib.native import _find_theme


class _Slide:
    """Handle d'une slide clonée, pour réécrire son texte par ancre."""

    def __init__(self, slide):
        self.slide = slide

    def set(self, swaps: dict[str, str]):
        """swaps = {texte_actuel (préfixe): nouveau_texte}. Réécrit la 1re forme qui matche."""
        for cur, new in swaps.items():
            sh = self._find(cur)
            if sh is None:
                print(f"[ModelDeck] ancre introuvable, ignorée : « {cur[:40]} »")
                continue
            _set_text(sh, new)
        return self

    def _find(self, current: str):
        for sh in self.slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(current):
                return sh
        return None


class ModelDeck:
    def __init__(self, theme: str):
        self.dir = _find_theme(theme)
        cat = json.loads((self.dir / "catalog.json").read_text("utf-8"))
        if cat.get("mode") != "models":
            raise ValueError(
                f"thème « {theme} » est en mode '{cat.get('mode')}' (masques) → "
                "utiliser NativeDeck, pas ModelDeck.")
        self.models = cat["models"]
        self.prs = Presentation(str(self.dir / "template.pptx"))
        self._sources = list(self.prs.slides)          # slides d'exemple = sources
        self._n_sources = len(self._sources)

    def catalog(self) -> list[dict]:
        return self.models

    def _resolve(self, ref):
        """ref = index (1-based) | nom | kind → slide source."""
        if isinstance(ref, int):
            return self._sources[ref - 1]
        for m in self.models:
            if m["name"] == ref or m["kind"] == ref:
                return self._sources[m["index"] - 1]
        # repli : match partiel sur le nom
        for m in self.models:
            if ref.lower() in m["name"].lower():
                return self._sources[m["index"] - 1]
        raise KeyError(f"modèle introuvable : {ref!r}")

    def clone(self, ref) -> _Slide:
        """Duplique la slide-modèle `ref` à la fin du deck et renvoie un handle."""
        src = self._resolve(ref)
        new = self.prs.slides.add_slide(src.slide_layout)
        for shp in list(new.shapes):                   # vider les placeholders ajoutés
            shp._element.getparent().remove(shp._element)
        for shp in src.shapes:                          # recopier les formes designées
            new.shapes._spTree.append(copy.deepcopy(shp._element))
        _remap_rels(src, new)                           # réparer les images (r:embed)
        return _Slide(new)

    def save(self, path: str) -> str:
        from pathlib import Path
        # retirer les slides-sources : ne garder que les clones
        lst = self.prs.slides._sldIdLst
        for sldId in list(lst)[:self._n_sources]:
            self.prs.part.drop_rel(sldId.get(qn("r:id")))
            lst.remove(sldId)
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return str(path)

    def export_pdf(self, pptx_path: str, pdf_path: str | None = None):
        from lib.pdf_export import pptx_to_pdf, soffice_path
        if soffice_path() is None:
            return None
        return pptx_to_pdf(pptx_path, pdf_path)


def _set_text(shape, text: str):
    """Remplace le texte d'une forme en gardant le formatage de son 1er run."""
    if not shape.has_text_frame:
        return
    p = shape.text_frame.paragraphs[0]
    if not p.runs:
        p.text = text
        return
    p.runs[0].text = text
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)


def _remap_rels(src_slide, new_slide):
    """Recopie les parts liées (images…) du modèle vers le clone et remappe les rId."""
    src_rels = src_slide.part.rels
    for el in new_slide.shapes._spTree.iter():
        for attr in (qn("r:embed"), qn("r:link")):
            rId = el.get(attr)
            if not rId or rId not in src_rels:
                continue
            rel = src_rels[rId]
            if rel.is_external:
                new_id = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_id = new_slide.part.relate_to(rel.target_part, rel.reltype)
            el.set(attr, new_id)
