"""NativeTemplate — ingère un template PPTX client EN ENTIER comme thème.

Principe (cf. évolution slider) : un thème n'est pas un nuancier réduit. C'est
le template du client conservé tel quel (master + slideLayouts + assets), plus :
  • un CATALOGUE des layouts (placeholders avec rôle déduit + signature) qui
    permet ensuite de remplir n'importe quel masque par RÔLE (pas par idx) ;
  • des TOKENS dérivés du thème PPTX (clrScheme + fontScheme), pour recolorer
    icônes / cadrer images, en métadonnée — le look vient du template, pas d'eux.

Sortie de `ingest()` dans <out>/<name>/ :
    template.pptx     copie fidèle du template (le système de design, intact)
    catalog.json      layouts → placeholders (idx, type, rôle, géométrie) + signature
    tokens.json       palette + polices dérivées du thème
    assets/photo/     photos du template (jpeg/png, + wdp converti si possible)
"""
from __future__ import annotations

import colorsys
import json
import re
import shutil
import subprocess
import zipfile
from pathlib import Path

try:  # parsing XML durci si dispo ; repli stdlib (entrée = PPTX local de confiance)
    from defusedxml.ElementTree import fromstring as _xml_fromstring
except ImportError:
    from xml.etree.ElementTree import fromstring as _xml_fromstring

from pptx import Presentation

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
EMU_CM = 360000.0

_WEIGHT_WORDS = ("Thin", "ExtraLight", "Light", "Regular", "Medium", "SemiBold",
                 "Semibold", "Bold", "ExtraBold", "Extrabold", "Black", "Display",
                 "Heavy", "Book", "Italic")


def _cm(emu) -> float:
    return round((emu or 0) / EMU_CM, 2)


def _base_family(name: str) -> str:
    """« Urbanist Black_GEL » → « Urbanist » ; « Aptos Display » → « Aptos »."""
    if not name:
        return "Inter"
    name = name.split("_")[0].strip()
    words = name.split()
    kept = [w for w in words if w not in _WEIGHT_WORDS]
    return " ".join(kept) if kept else words[0]


def _sat(hex6: str) -> float:
    r, g, b = (int(hex6[i:i + 2], 16) / 255 for i in (0, 2, 4))
    return colorsys.rgb_to_hls(r, g, b)[2]


def _lum(hex6: str) -> float:
    r, g, b = (int(hex6[i:i + 2], 16) / 255 for i in (0, 2, 4))
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


# ---------------------------------------------------------------- thème → tokens

def theme_scheme(pptx_path: str) -> tuple[dict, str]:
    """Lit le premier theme*.xml : clrScheme (dk1/lt1/dk2/lt2/accent1-6) + police majeure."""
    with zipfile.ZipFile(pptx_path) as z:
        themes = sorted(n for n in z.namelist()
                        if re.match(r"ppt/theme/theme\d+\.xml$", n))
        xml = z.read(themes[0]).decode("utf-8")
    root = _xml_fromstring(xml)
    scheme: dict[str, str] = {}
    clr = root.find(f".//{A}clrScheme")
    for child in clr or []:
        tag = child.tag.split("}")[1]
        srgb = child.find(f"{A}srgbClr")
        sysc = child.find(f"{A}sysClr")
        hexv = (srgb.get("val") if srgb is not None
                else sysc.get("lastClr") if sysc is not None else None)
        if hexv:
            scheme[tag] = hexv.upper()
    major = root.find(f".//{A}fontScheme/{A}majorFont/{A}latin")
    family = _base_family(major.get("typeface") if major is not None else "")
    return scheme, family


def build_tokens(name: str, scheme: dict, family: str, source: str) -> dict:
    """Mappe le clrScheme du template sur les tokens slider (heuristique, raffinable)."""
    def hx(k, default):
        return "#" + scheme.get(k, default)

    accents = [scheme[k] for k in ("accent1", "accent2", "accent3", "accent4",
                                   "accent5", "accent6") if k in scheme]
    # signature = accent le plus saturé ; primary = couleur de marque lisible
    signature = max(accents, key=_sat) if accents else scheme.get("accent1", "2563EB")
    dk2 = scheme.get("dk2")
    primary = dk2 if (dk2 and _sat(dk2) > 0.25) else (
        max((a for a in accents if 0.15 < _lum(a) < 0.7), key=_sat, default=signature))

    return {
        "$schema": "../_schema/charte.schema.json",
        "name": name,
        "label": name,
        "source": source,
        "notes": "Tokens DÉRIVÉS du thème PPTX (clrScheme/fontScheme). Le look provient "
                 "du template préservé (template.pptx + catalog.json), pas de ces tokens — "
                 "ils servent à recolorer icônes / cadrer images. Rôles raffinables.",
        "colors": {
            "primary":      {"value": "#" + primary,   "name": "Marque",      "usage": "accent texte lisible"},
            "primary-deep": {"value": hx("dk2", "08212C"), "name": "Marque foncée", "usage": "fonds sombres, dividers"},
            "signature":    {"value": "#" + signature, "name": "Signature",   "usage": "barres, carrés, marqueurs"},
            "text":         {"value": hx("dk1", "1A1A1A"), "name": "Encre",    "usage": "texte courant et titres"},
            "text-strong":  {"value": hx("dk1", "000000"), "name": "Encre forte", "usage": "titres contrastés"},
            "muted":        {"value": "#5D848A",       "name": "Gris",        "usage": "secondaire"},
            "rule":         {"value": "#D5E0E2",       "name": "Filet",       "usage": "bordures"},
            "bg":           {"value": hx("lt1", "FFFFFF"), "name": "Fond",     "usage": "fond par défaut"},
            "bg-soft":      {"value": hx("lt2", "EFF4F5"), "name": "Fond doux", "usage": "panneaux"},
        },
        "_palette_brute": {k: "#" + v for k, v in scheme.items()},
        "fonts": {
            "primary": {
                "family": family,
                "fallback": ["Inter", "Helvetica Neue", "Arial", "sans-serif"],
                "weights": [300, 400, 500, 600, 700, 800, 900],
            }
        },
    }


# ---------------------------------------------------------------- catalogue

_PH_TYPE = {
    "TITLE (1)": "title", "CENTER_TITLE (3)": "title", "SUBTITLE (4)": "subtitle",
    "SLIDE_NUMBER (13)": "page_number", "FOOTER (15)": "footer", "DATE (16)": "date",
    "PICTURE (18)": "picture",
}


def _role(ph, title_top, title_left) -> str:
    """Rôle déduit par TYPE puis POSITION (jamais par idx)."""
    t = str(ph.placeholder_format.type)
    if t in _PH_TYPE:
        return _PH_TYPE[t]
    # BODY (2) et apparentés : position relative au titre
    top, left = ph.top or 0, ph.left or 0
    w = ph.width or 0
    if w < 1.2 * EMU_CM:
        return "decoration"                      # barres / filets verticaux décoratifs
    # (les footnotes sont fines en hauteur mais larges → pas de garde sur h)
    if top > 16.5 * EMU_CM:
        return "footnote"
    if title_top is not None:
        # eyebrow = juste AU-DESSUS du titre ET aligné horizontalement avec lui
        # (sinon un bloc de contenu haut placé à droite serait pris pour un eyebrow
        #  quand le titre est bas dans la slide — cas « Title + 3 Blocs »).
        near_above = title_top - 2.5 * EMU_CM <= top < title_top - 0.3 * EMU_CM
        aligned = abs(left - (title_left or 0)) < 3 * EMU_CM
        if near_above and aligned:
            return "eyebrow"
        if title_top <= top <= title_top + 2.2 * EMU_CM and w > 18 * EMU_CM:
            return "subtitle"
    return "content"


def catalog(pptx_path: str) -> list[dict]:
    prs = Presentation(pptx_path)
    layouts: list[dict] = []
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            title = next((p for p in lay.placeholders
                          if str(p.placeholder_format.type) in ("TITLE (1)", "CENTER_TITLE (3)")), None)
            title_top = title.top if title is not None else None
            title_left = title.left if title is not None else None
            phs = []
            for ph in lay.placeholders:
                f = ph.placeholder_format
                phs.append({
                    "idx": f.idx, "type": str(f.type), "role": _role(ph, title_top, title_left),
                    "x": _cm(ph.left), "y": _cm(ph.top), "w": _cm(ph.width), "h": _cm(ph.height),
                })
            roles = [p["role"] for p in phs]
            sig = {
                "content": roles.count("content"),
                "pictures": roles.count("picture"),
                "has_title": "title" in roles,
                "has_eyebrow": "eyebrow" in roles,
                "has_subtitle": "subtitle" in roles,
            }
            layouts.append({"name": lay.name, "kind": _kind(lay.name, sig, phs),
                            "signature": sig, "placeholders": phs})
    return layouts


def _kind(name: str, sig: dict, phs: list[dict]) -> str:
    """Étiquette indicative (le binder pourra l'affiner). Signature d'abord, nom en repli."""
    big_pic = any(p["type"] == "PICTURE (18)" and p["w"] * p["h"] > 250 for p in phs)
    n = sig["content"]
    low = name.lower()
    if "cover" in low:
        return "cover"
    if "end slide" in low or "contact" in low:
        return "end"
    if "agenda" in low:
        return "agenda"
    if big_pic and n <= 1:
        return "section"          # chapitre / divider photo
    if "highlight" in low:
        return "statement"
    if n >= 2:
        return f"grid-{n}"
    if n == 1:
        return "list-or-text"
    return "content"


# ------------------------------------------------ slides-modèles (mode « models »)
# Pour les templates SANS bibliothèque de masques : le design est dessiné sur les
# slides elles-mêmes. On catalogue chaque slide d'exemple comme MODÈLE réutilisable
# (à cloner + réécrire par ancre de texte — cf. lib/models.ModelDeck).

def _anchor_role(text: str, x: float, y: float, w: float, H: float) -> str:
    low = text.strip()
    if y > 0.85 * H:
        if len(low) <= 4 and any(c.isdigit() for c in low):
            return "page_number"
        return "footer"
    if len(low) > 45:
        return "body"
    if low.isupper() and len(low) <= 30:
        return "section"          # tag de section / eyebrow en capitales
    if y < 0.45 * H:
        return "title"
    return "label"


def slide_models(pptx_path: str) -> list[dict]:
    prs = Presentation(pptx_path)
    H = (prs.slide_height or 1) / EMU_CM
    out: list[dict] = []
    for i, s in enumerate(prs.slides, 1):
        anchors, has_image = [], False
        for sh in s.shapes:
            if str(sh.shape_type).startswith("PICTURE"):
                has_image = True
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if not t:
                continue
            anchors.append({
                "shape_id": sh.shape_id, "text": t.replace("\n", " ")[:90],
                "role": _anchor_role(t, _cm(sh.left), _cm(sh.top), _cm(sh.width), H),
                "x": _cm(sh.left), "y": _cm(sh.top), "w": _cm(sh.width),
            })
        title = next((a["text"] for a in anchors if a["role"] == "title"), None)
        out.append({"index": i, "name": title or f"slide {i}",
                    "kind": _model_kind(anchors), "has_image": has_image,
                    "anchors": anchors})
    return out


def _model_kind(anchors: list[dict]) -> str:
    txts = " ".join(a["text"].lower() for a in anchors)
    n_body = sum(1 for a in anchors if a["role"] == "body")
    if "/" in txts and any(a["role"] == "section" for a in anchors) and n_body == 0:
        return "divider"          # ex. « 01 / 02 » + tag de partie
    if len(anchors) <= 3:
        return "cover"            # peu de texte = couverture / clôture
    return "content"


def _has_layout_library(layouts_cat: list[dict]) -> bool:
    """Vrai si le template porte une vraie bibliothèque de masques (mode layouts)."""
    real = [l for l in layouts_cat
            if l["name"].strip().upper() not in ("DEFAULT", "BLANK", "VIDE")
            and (l["signature"]["content"] > 0 or l["signature"]["pictures"] > 0)]
    return len(real) >= 3


# ---------------------------------------------------------------- assets

def harvest_photos(pptx_path: str, dest: Path) -> list[str]:
    """Copie les photos du template (jpeg/jpg/png « grandes »), convertit le wdp si possible."""
    dest.mkdir(parents=True, exist_ok=True)
    out: list[str] = []
    has_convert = shutil.which("convert") is not None
    with zipfile.ZipFile(pptx_path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        for n in media:
            ext = n.rsplit(".", 1)[-1].lower()
            base = Path(n).name
            data = z.read(n)
            if ext in ("jpeg", "jpg", "png") and len(data) > 40_000:
                (dest / base).write_bytes(data)
                out.append(base)
            elif ext == "wdp" and has_convert:
                src = dest / base
                src.write_bytes(data)
                png = dest / (Path(base).stem + ".png")
                try:
                    subprocess.run(["convert", str(src), str(png)], check=True,
                                   capture_output=True, timeout=30)
                    out.append(png.name)
                finally:
                    src.unlink(missing_ok=True)
    return out


# ---------------------------------------------------------------- ingest

def ingest(src_pptx: str, name: str, out_dir: str) -> dict:
    src = Path(src_pptx)
    if not src.exists():
        raise FileNotFoundError(src)
    theme = Path(out_dir) / name
    (theme / "assets").mkdir(parents=True, exist_ok=True)

    # 1. préserver le template tel quel
    shutil.copy(src, theme / "template.pptx")

    # 2. catalogue — masques (mode layouts) ET/OU slides-modèles (mode models)
    cat = catalog(src_pptx)
    models = slide_models(src_pptx)
    mode = "layouts" if _has_layout_library(cat) else "models"
    (theme / "catalog.json").write_text(
        json.dumps({"source": src.name, "mode": mode,
                    "layouts": cat, "models": models}, ensure_ascii=False, indent=2),
        encoding="utf-8")

    # 3. tokens dérivés du thème
    scheme, family = theme_scheme(src_pptx)
    tokens = build_tokens(name, scheme, family, f"ingest de {src.name} (template préservé)")
    (theme / "tokens.json").write_text(
        json.dumps(tokens, ensure_ascii=False, indent=2), encoding="utf-8")

    # 4. photos
    photos = harvest_photos(src_pptx, theme / "assets" / "photo")

    return {"theme": str(theme), "mode": mode, "layouts": len(cat),
            "models": len(models), "family": family,
            "palette": len(scheme), "photos": len(photos)}
